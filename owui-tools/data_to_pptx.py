"""
title: Data to PPTX
author: corporate
version: 0.1.0
requirements: python-pptx
description: Генерирует PPTX из структурированных данных. Поддерживает типы слайдов: title, bullets, table. Использует корпоративный шаблон, если он есть.
"""

import io
import base64
import json
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt


class Tools:
    def __init__(self):
        self.template_path = "/app/backend/data/templates/corp_template.pptx"

    def data_to_pptx(
        self,
        slides_json: str,
        output_filename: Optional[str] = "presentation.pptx",
        use_template: bool = True,
    ) -> dict:
        """
        Сгенерировать PPTX из структурированных данных.

        slides_json — JSON-массив слайдов вида:
        [
          {"type": "title", "title": "Заголовок", "subtitle": "Подзаголовок"},
          {"type": "bullets", "title": "Итоги", "bullets": ["Пункт 1", "Пункт 2"]},
          {"type": "table", "title": "Метрики",
           "headers": ["Метрика", "Q1", "Q2"],
           "rows": [["Выручка", "100", "120"], ["Маржа", "20%", "22%"]]}
        ]
        """
        slides = json.loads(slides_json)

        try:
            prs = Presentation(self.template_path) if use_template else Presentation()
        except Exception:
            prs = Presentation()

        for slide_def in slides:
            stype = slide_def.get("type", "bullets")

            if stype == "title":
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_def.get("title", "")
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_def.get("subtitle", "")

            elif stype == "bullets":
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_def.get("title", "")
                body = slide.placeholders[1].text_frame
                body.clear()
                for i, b in enumerate(slide_def.get("bullets", [])):
                    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    p.text = b
                    p.font.size = Pt(20)

            elif stype == "table":
                layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = slide_def.get("title", "")
                headers = slide_def.get("headers", [])
                rows = slide_def.get("rows", [])
                if headers and rows:
                    n_rows, n_cols = len(rows) + 1, len(headers)
                    left, top = Inches(0.5), Inches(1.5)
                    width, height = Inches(9), Inches(0.4 * n_rows)
                    tbl = slide.shapes.add_table(
                        n_rows, n_cols, left, top, width, height
                    ).table
                    for c, h in enumerate(headers):
                        tbl.cell(0, c).text = str(h)
                    for r, row in enumerate(rows, start=1):
                        for c, val in enumerate(row):
                            tbl.cell(r, c).text = str(val)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return {
            "filename": output_filename,
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "content_base64": base64.b64encode(out.read()).decode(),
            "slides_count": len(slides),
        }
