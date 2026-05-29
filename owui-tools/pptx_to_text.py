"""
title: PPTX to Text
author: corporate
version: 0.1.0
requirements: python-pptx
description: Извлекает весь текст из PPTX — заголовки, тело, таблицы, заметки докладчика. Используется моделью для чтения/саммари презентаций.
"""

import io
import base64
from typing import Dict, Any
from pptx import Presentation


class Tools:
    def __init__(self):
        pass

    def pptx_to_text(self, pptx_base64: str) -> Dict[str, Any]:
        """
        Извлечь весь текст из PPTX: заголовки, тело, таблицы, заметки.
        """
        data = base64.b64decode(pptx_base64)
        prs = Presentation(io.BytesIO(data))

        slides_out = []
        for i, slide in enumerate(prs.slides, start=1):
            texts, tables = [], []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in p.runs).strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([c.text.strip() for c in row.cells])
                    tables.append(rows)
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides_out.append({
                "slide": i,
                "text": texts,
                "tables": tables,
                "notes": notes,
            })
        return {"slides_count": len(slides_out), "slides": slides_out}
